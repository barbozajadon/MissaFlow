import argparse
import re

from pptx import Presentation

from database.database import get_session, init_db
from database.models import Hymn
from utils.resources import resource_path

# -------------------------------------------------
# Configuration
# -------------------------------------------------
PPT_PATH = resource_path(
    "assets",
    "Joyful Celebration Hymn Book.pptx"
)

# Detects:
# 39 - Amazing Grace
# 102 - Come Back To Me
TITLE_PATTERN = re.compile(r"^\s*(\d+)\s*-\s*(.+)$")


# -------------------------------------------------
# Read a slide's text, separating out the header line (if present)
# from the rest of the lyric lines.
#
# IMPORTANT: the header ("15 - A Lamp For Our Steps") is NOT reliably
# the first line of text on its slide - it lives in its own text box,
# and shape order on the slide doesn't match visual/reading order (the
# lyrics text box can come before the title box in the shape list).
# So every shape's every line is checked for the header pattern, not
# just the first line of the slide.
# -------------------------------------------------
def find_title_and_lyrics(slide):
    title_match = None
    lyrics_lines = []
    for shape in slide.shapes:
        if not hasattr(shape, "text") or not shape.text.strip():
            continue
        for line in shape.text.split("\n"):
            line = line.strip()
            if not line:
                continue
            if title_match is None:
                match = TITLE_PATTERN.match(line)
                if match:
                    title_match = match
                    continue  # the header line itself isn't lyrics
            lyrics_lines.append(line)
    return title_match, lyrics_lines


# -------------------------------------------------
# Main Import
# -------------------------------------------------
def import_hymns(dry_run: bool = False):
    if not dry_run:
        init_db()

    prs = Presentation(PPT_PATH)
    hymns = []
    current = None

    for index, slide in enumerate(prs.slides):
        slide_number = index + 1
        title_match, lines = find_title_and_lyrics(slide)

        if title_match:
            hymn_number = title_match.group(1).strip()
            title = title_match.group(2).strip()

            if current and hymn_number == current["number"]:
                # Same hymn number as before - this is a repeated
                # refrain/response section within the SAME hymn (some
                # hymns re-show their header at each repeat), not a new
                # hymn. Extend the current entry instead of splitting it.
                current["end_slide"] = slide_number
                current["lyrics_lines"].extend(lines)
                continue

            # Finish previous hymn before starting this new one
            if current:
                current["end_slide"] = slide_number - 1
                hymns.append(current)

            current = {
                "number": hymn_number,
                "title": title,
                "start_slide": slide_number,
                "end_slide": slide_number,
                "lyrics_lines": lines,
            }
        else:
            # Continuation slide (no header) - more lyrics for the
            # current hymn
            if current:
                current["lyrics_lines"].extend(lines)
                current["end_slide"] = slide_number
            else:
                print(f"[WARN] Slide {slide_number} has no title and no current "
                      f"hymn context - skipping.")

    # Finish last hymn
    if current:
        current["end_slide"] = len(prs.slides)
        hymns.append(current)

    for hymn in hymns:
        hymn["lyrics_text"] = "\n".join(hymn["lyrics_lines"]).strip()

    # ----------------------------------------------------------------
    # Dry run: print everything, write nothing
    # ----------------------------------------------------------------
    if dry_run:
        print(f"Found {len(hymns)} hymn(s) in {PPT_PATH}\n")
        for hymn in hymns:
            preview = hymn["lyrics_text"][:60].replace("\n", " / ")
            print(
                f"  #{hymn['number']:<6} {hymn['title']:<35} "
                f"slides {hymn['start_slide']}-{hymn['end_slide']:<5} "
                f"lyrics: {preview}..."
            )
        print("\nDry run only - nothing written to the database.")
        return

    # ----------------------------------------------------------------
    # Save into database
    # ----------------------------------------------------------------
    with get_session() as session:
        existing_numbers = {h.hymn_number for h in session.query(Hymn.hymn_number).all()}
        inserted = 0
        skipped = 0
        for hymn in hymns:
            if hymn["number"] in existing_numbers:
                skipped += 1
                continue
            session.add(
                Hymn(
                    hymn_number=hymn["number"],
                    title=hymn["title"],
                    lyrics=hymn["lyrics_text"] or None,
                    start_slide=hymn["start_slide"],
                    end_slide=hymn["end_slide"],
                    category=None,
                    language="English",
                )
            )
            existing_numbers.add(hymn["number"])  # so later duplicates in this same run are caught too
            inserted += 1
        session.commit()

    print(f"Imported {inserted} new hymn(s), skipped {skipped} already-existing, out of {len(hymns)} found.")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Import hymns from the master PPTX into the database.")
    parser.add_argument(
        "--dry-run",
        action="store_true",
        help="Preview extracted hymns without writing to the database.",
    )
    args = parser.parse_args()

    import_hymns(dry_run=args.dry_run)
