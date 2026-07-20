"""
import_psalm_responses.py

Imports a monthly Psalm Response PPTX (one slide per calendar date, e.g.
"7th June, Sunday" + that day's response text) into the database:

1. Each slide becomes a Hymn-table row (category='psalm_response') so it
   can be found via the normal Search dialog and has a real slide
   reference for presenting.
2. Each date's LiturgicalCalendar entry is updated with the response
   TEXT (for the Planner's auto-fill label) AND a link to that Hymn row
   (psalm_response_hymn_id), so the Planner can auto-select the actual
   slide for that date's "Psalm Response" slot.

USAGE
-----
    pip install python-pptx

    # Preview without writing to the database
    python import_psalm_responses.py "assets/June_Psalm_Responses.pptx" --year 2026 --dry-run

    # Import for real
    python import_psalm_responses.py "assets/June_Psalm_Responses.pptx" --year 2026

Run this once per month, pointing at that month's file - re-running is
safe, existing rows for the same (source file, slide number) are reused
rather than duplicated.
"""
from __future__ import annotations

import argparse
import datetime
import re

from pptx import Presentation

from database.database import get_session, init_db
from database.models import Hymn
from services import calendar_service

# Matches "1st June, Monday", "21st June, Sunday", etc.
DATE_LINE_PATTERN = re.compile(
    r"^\s*(\d{1,2})(?:st|nd|rd|th)?\s+([A-Za-z]+)\s*,?\s*[A-Za-z]*\s*$"
)


def get_slide_lines(slide) -> list[str]:
    lines = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            for line in shape.text.split("\n"):
                line = line.strip()
                if line:
                    lines.append(line)
    return lines


def parse_date_line(date_line: str, year: int) -> datetime.date:
    """Parse '1st June, Monday' (day name ignored) into a real date."""
    match = DATE_LINE_PATTERN.match(date_line)
    if not match:
        raise ValueError(f"Could not parse date from: {date_line!r}")
    day = int(match.group(1))
    month_name = match.group(2)
    return datetime.datetime.strptime(f"{day} {month_name} {year}", "%d %B %Y").date()


def extract_entries(ppt_path: str, year: int) -> list[dict]:
    prs = Presentation(ppt_path)
    entries = []
    for index, slide in enumerate(prs.slides):
        slide_number = index + 1
        lines = get_slide_lines(slide)
        if len(lines) < 2:
            print(f"[WARN] Slide {slide_number} has fewer than 2 text lines, skipping: {lines}")
            continue

        date_line, response_line = lines[0], lines[1]
        try:
            mass_date = parse_date_line(date_line, year)
        except ValueError as exc:
            print(f"[WARN] Slide {slide_number}: {exc} - skipping")
            continue

        entries.append({
            "date": mass_date,
            "response_text": response_line,
            "slide_number": slide_number,
            "date_line": date_line,
        })
    return entries


def import_psalm_responses(ppt_path: str, year: int, dry_run: bool = False) -> None:
    entries = extract_entries(ppt_path, year)

    print(f"Found {len(entries)} dated psalm response(s) in {ppt_path}\n")

    if dry_run:
        for e in entries:
            print(f"  {e['date'].isoformat()} (slide {e['slide_number']:<4}) "
                  f"{e['response_text'][:60]}")
        print("\nDry run only - nothing written to the database.")
        return

    init_db()

    # --- Phase 1: create/reuse Hymn rows for each slide, commit, collect IDs ---
    hymn_ids_by_slide: dict[int, int] = {}
    with get_session() as session:
        inserted = 0
        for e in entries:
            existing = session.query(Hymn).filter_by(
                category="psalm_response",
                source_file=ppt_path,
                start_slide=e["slide_number"],
            ).first()
            if existing:
                hymn_ids_by_slide[e["slide_number"]] = existing.id
                continue

            hymn = Hymn(
                hymn_number=None,
                title=f"Psalm Response - {e['date'].isoformat()}",
                lyrics=e["response_text"],
                category="psalm_response",
                language="English",
                source_file=ppt_path,
                start_slide=e["slide_number"],
                end_slide=e["slide_number"],
            )
            session.add(hymn)
            session.flush()  # populate hymn.id before commit
            hymn_ids_by_slide[e["slide_number"]] = hymn.id
            inserted += 1
        session.commit()

    print(f"Created {inserted} new psalm-response slide entr(y/ies) "
          f"(reused {len(entries) - inserted} existing).")

    # --- Phase 2: upsert each date's LiturgicalCalendar entry, linking the Hymn ---
    linked = 0
    for e in entries:
        hymn_id = hymn_ids_by_slide.get(e["slide_number"])
        calendar_service.upsert_entry(
            mass_date=e["date"],
            psalm_response=e["response_text"],
            psalm_response_hymn_id=hymn_id,
        )
        linked += 1

    print(f"Linked {linked} calendar date(s) to their psalm-response slide.")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="Import a monthly Psalm Response PPTX into the database."
    )
    parser.add_argument("ppt_path", help="Path to the monthly psalm response .pptx file")
    parser.add_argument(
        "--year", type=int, required=True,
        help="Calendar year this file's dates belong to (the PPT only has day+month, no year)",
    )
    parser.add_argument(
        "--dry-run", action="store_true",
        help="Preview extracted entries without writing to the database",
    )
    args = parser.parse_args()

    import_psalm_responses(args.ppt_path, args.year, dry_run=args.dry_run)
