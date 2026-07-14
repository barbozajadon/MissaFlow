"""
PowerPoint generation for a planned Mass.

Copies the slide ranges of each selected hymn (and any mass-part slides,
if you keep those in the master deck too) out of master_hymnal.pptx and
assembles them, in Mass order, into a new presentation saved as
YYYY-MM-DD_Mass.pptx inside generated_presentations/.

FIDELITY NOTE
-------------
python-pptx has no built-in "copy slide" operation. This module copies
each shape's XML element (text, formatting, tables) and re-adds pictures
via their original image bytes, which preserves layout, text formatting,
and images for the vast majority of hymn slides. Slide transitions and
animations live in slide-level XML that is not exposed by python-pptx's
object model, so on a deck with custom transitions those may not
carry over perfectly - if that matters for your Mass slides, verify the
first generated file after your first real run.
"""
from __future__ import annotations

import copy
import datetime
import logging
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.slide import Slide

logger = logging.getLogger(__name__)


def _copy_shapes(source_slide: Slide, dest_slide: Slide) -> None:
    """Deep-copy every shape's XML element from source to dest slide."""
    for shape in source_slide.shapes:
        if shape.shape_type is not None and shape.shape_type == 13:  # PICTURE
            _copy_picture(shape, source_slide, dest_slide)
        else:
            new_el = copy.deepcopy(shape._element)
            dest_slide.shapes._spTree.append(new_el)


def _copy_picture(shape, source_slide: Slide, dest_slide: Slide) -> None:
    """Re-add a picture shape using the original image bytes and position."""
    try:
        image_part = shape.image
        image_bytes = image_part.blob
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
        import io
        dest_slide.shapes.add_picture(io.BytesIO(image_bytes), left, top, width, height)
    except Exception:
        logger.exception("Failed to copy picture shape, falling back to raw XML copy")
        new_el = copy.deepcopy(shape._element)
        dest_slide.shapes._spTree.append(new_el)


def _copy_background(source_slide: Slide, dest_slide: Slide) -> None:
    """Copy a slide-level background override, if the source slide has one."""
    bg = source_slide.element.find(
        "{http://schemas.openxmlformats.org/presentationml/2006/main}bg"
    )
    if bg is not None:
        dest_slide.element.insert(0, copy.deepcopy(bg))


def copy_slide(prs: Presentation, source_slide: Slide) -> Slide:
    """Append a copy of source_slide (using its own layout) to prs."""
    layout = source_slide.slide_layout
    dest_slide = prs.slides.add_slide(layout)

    # Remove placeholder shapes python-pptx auto-added from the layout,
    # so we don't end up with duplicate/empty placeholders under our copy.
    for shape in list(dest_slide.shapes):
        shape._element.getparent().remove(shape._element)

    _copy_background(source_slide, dest_slide)
    _copy_shapes(source_slide, dest_slide)
    return dest_slide


def generate_presentation(
    master_pptx_path: str,
    slide_ranges: list[tuple[int, int]],
    mass_date: datetime.date,
    output_dir: str,
) -> str:
    """
    Build a new presentation containing the given 1-indexed slide ranges,
    pulled in order from master_pptx_path, and save it to output_dir.

    slide_ranges: list of (start_slide, end_slide) tuples, INCLUSIVE,
        1-indexed, in the exact order they should appear - i.e. build
        this list from the ordered MassItems of the MassPlan before
        calling this function.

    Returns the path to the generated .pptx file.
    """
    master_path = Path(master_pptx_path)
    if not master_path.exists():
        raise FileNotFoundError(f"Master hymnal not found: {master_pptx_path}")

    source_prs = Presentation(str(master_path))
    source_slides = list(source_prs.slides)
    total_slides = len(source_slides)

    # Build the new deck from the master's own template so the theme/
    # slide-size/master matches exactly, then strip the default slide.
    out_prs = Presentation(str(master_path))
    while len(out_prs.slides._sldIdLst) > 0:
        out_prs.slides._sldIdLst.remove(out_prs.slides._sldIdLst[0])

    copied_count = 0
    for start, end in slide_ranges:
        if start < 1 or end > total_slides or start > end:
            logger.warning("Skipping invalid slide range (%s, %s); deck has %s slides",
                            start, end, total_slides)
            continue
        for slide_no in range(start, end + 1):
            source_slide = source_slides[slide_no - 1]
            copy_slide(out_prs, source_slide)
            copied_count += 1

    if copied_count == 0:
        raise ValueError("No valid slide ranges were provided - nothing to generate")

    output_path = Path(output_dir)
    output_path.mkdir(parents=True, exist_ok=True)
    filename = f"{mass_date.isoformat()}_Mass.pptx"
    full_path = output_path / filename

    out_prs.save(str(full_path))
    logger.info("Generated presentation with %s slides at %s", copied_count, full_path)
    return str(full_path)


def build_slide_ranges_from_mass_items(mass_items: list) -> list[tuple[int, int]]:
    """
    Convert an ordered list of MassItem ORM objects (each with a .hymn
    relationship carrying start_slide/end_slide) into the slide_ranges
    list generate_presentation() expects. Items without a linked hymn
    (e.g. fixed spoken parts with only text_override) are skipped, since
    they have no slides in the master deck.
    """
    ranges: list[tuple[int, int]] = []
    for item in mass_items:
        hymn = getattr(item, "hymn", None)
        if hymn and hymn.start_slide and hymn.end_slide:
            ranges.append((hymn.start_slide, hymn.end_slide))
    return ranges
