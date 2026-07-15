import argparse
import re

from pptx import Presentation

from database.database import get_session, init_db
from database.models import Hymn

# -------------------------------------------------
# Configuration
# -------------------------------------------------
PPT_PATH = "assets/Joyful Celebration Hymn Book.pptx"

# Detects:
# 39 - Amazing Grace
# 102 - Come Back To Me
TITLE_PATTERN = re.compile(r"^\s*(\d+)\s*-\s*(.+)$")


# -------------------------------------------------
# Read all text from a slide, line by line
# -------------------------------------------------
def get_slide_lines(slide):
    lines = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            for line in shape.text.split("\n"):
                line = line.strip()
                if line:
                    lines.append(line)
    return lines


# -------------------------------------------------
# Main Import
# -------------------------------------------------
def import_hymns(dry_run: bool = False):
    if not dry_run:
        init_db()

    prs = Presentation(PPT_PATH)
    hymns = []
    current = None

    for index, slide in enumerate(prs.slides):
        slide_number = index + 1
        lines = get_slide_lines(slide)
        first_line = lines[0] if lines else ""
        match = TITLE_PATTERN.match(first_line)

        if match:
            hymn_number = match.group(1).strip()
            title = match.group(2).strip()
            remaining_lines = lines[1:]  # rest of this slide's text is lyrics too

            if current and hymn_number == current["number"]:
                # Same hymn number as the slide before - this is the header
                # repeating on a later slide of the SAME hymn (e.g. every
                # verse's slide re-shows "141 - Title"), not a new hymn.
                # Extend the current hymn instead of starting a new one.
                current["end_slide"] = slide_number
                current["lyrics_lines"].extend(remaining_lines)
                continue

            # Finish previous hymn before starting this new one
            if current:
                current["end_slide"] = slide_number - 1
                hymns.append(current)

            current = {
                "number": hymn_number,
                "title": title,
                "start_slide": slide_number,
                "end_slide": slide_number,
                "lyrics_lines": remaining_lines,
            }
        else:
            # Continuation slide - append its text as more lyrics for the
            # current hymn
            if current:
                current["lyrics_lines"].extend(lines)

    # Finish last hymn
    if current:
        current["end_slide"] = len(prs.slides)
        hymns.append(current)

    for hymn in hymns:
        hymn["lyrics_text"] = "\n".join(hymn["lyrics_lines"]).strip()

    # ----------------------------------------------------------------
    # Dry run: print everything, write nothing
    # ----------------------------------------------------------------
    if dry_run:
        print(f"Found {len(hymns)} hymn(s) in {PPT_PATH}\n")
        for hymn in hymns:
            preview = hymn["lyrics_text"][:60].replace("\n", " / ")
            print(
                f"  #{hymn['number']:<6} {hymn['title']:<35} "
                f"slides {hymn['start_slide']}-{hymn['end_slide']:<5} "
                f"lyrics: {preview}..."
            )
        print("\nDry run only - nothing written to the database.")
        return

    # ----------------------------------------------------------------
    # Save into database
    # ----------------------------------------------------------------
    with get_session() as session:
        existing_numbers = {h.hymn_number for h in session.query(Hymn.hymn_number).all()}
        inserted = 0
        skipped = 0
        for hymn in hymns:
            if hymn["number"] in existing_numbers:
                skipped += 1
                continue
            session.add(
                Hymn(
                    hymn_number=hymn["number"],
                    title=hymn["title"],
                    lyrics=hymn["lyrics_text"] or None,
                    start_slide=hymn["start_slide"],
                    end_slide=hymn["end_slide"],
                    category=None,
                    language="English",
                )
            )
            existing_numbers.add(hymn["number"])  # so later duplicates in this same run are caught too
            inserted += 1
        session.commit()

    print(f"Imported {inserted} new hymn(s), skipped {skipped} already-existing, out of {len(hymns)} found.")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Import hymns from the master PPTX into the database.")
    parser.add_argument(
        "--dry-run",
        action="store_true",
        help="Preview extracted hymns without writing to the database.",
    )
    args = parser.parse_args()

    import_hymns(dry_run=args.dry_run)
